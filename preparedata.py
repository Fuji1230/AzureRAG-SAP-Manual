import json
import re
import os
import hashlib
import argparse
from base64 import urlsafe_b64encode
from openai import AzureOpenAI
from dotenv import load_dotenv
import requests
from dococr.parse_doc import get_content_from_document
from dococr.create_chunks import chunk_content
from azure.core.exceptions import ResourceNotFoundError
from azure.core.credentials import AzureKeyCredential
from azure.search.documents import SearchClient
from azure.search.documents.indexes import SearchIndexClient
from azure.cosmos import PartitionKey
from azure.cosmos.cosmos_client import CosmosClient
from azure.storage.blob import BlobServiceClient, ContainerClient
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from bs4 import BeautifulSoup
import urllib.parse

# 環境変数を読み込む
load_dotenv()

max_chunk_token_size = 2048
overlap_token_rate = 0
overlap_type = "NONE"

client = AzureOpenAI(
    azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
    api_key=os.getenv("AZURE_OPENAI_API_KEY"),
    api_version=os.getenv("AZURE_OPENAI_API_VERSION")
)

AZURE_OPENAI_CHAT_MODEL = os.getenv("AZURE_OPENAI_CHAT_MODEL")
AZURE_OPENAI_EMBED_MODEL = os.getenv("AZURE_OPENAI_EMBED_MODEL")
AZURE_OPENAI_CHAT_MAX_TOKENS = int(os.getenv("AZURE_OPENAI_CHAT_MAX_TOKENS", "1000"))

AI_SEARCH_ENDPOINT = os.getenv("AI_SEARCH_ENDPOINT")
AI_SEARCH_KEY = os.getenv("AI_SEARCH_KEY")
AI_SEARCH_API_VERSION = os.getenv("AI_SEARCH_API_VERSION", "2023-10-01-Preview")

credential = AzureKeyCredential(AI_SEARCH_KEY)
index_client = SearchIndexClient(
    endpoint=AI_SEARCH_ENDPOINT,
    credential=credential,
)

COSMOS_CONNECTION_STRING = os.getenv("COSMOS_CONNECTION_STRING")
COSMOS_DB_NAME = os.getenv("COSMOS_DB_NAME")
COSMOS_CONTAINER_NAME_KB = os.getenv("COSMOS_CONTAINER_NAME_KB")

cosmos_client = CosmosClient.from_connection_string(COSMOS_CONNECTION_STRING)
database = cosmos_client.get_database_client(COSMOS_DB_NAME)
database.create_container_if_not_exists(
    id=COSMOS_CONTAINER_NAME_KB,
    partition_key=PartitionKey(path=f"/id"),
)
container = database.get_container_client(COSMOS_CONTAINER_NAME_KB)

BLOB_STORAGE_CONNECTION_STRING = os.getenv("BLOB_STORAGE_CONNECTION_STRING")
BLOB_CONTAINER_NAME = os.getenv("BLOB_CONTAINER_NAME")
USE_BLOB_STORAGE = os.getenv("USE_BLOB_STORAGE")

blob_service_client = BlobServiceClient.from_connection_string(BLOB_STORAGE_CONNECTION_STRING)
blob_container_client = blob_service_client.get_container_client(BLOB_CONTAINER_NAME)

def extract_text_from_docx(path):
    doc = Document(path)
    return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

def extract_text_from_html(path):
    with open(path, 'r', encoding='utf-8') as f:
        html = f.read()
    soup = BeautifulSoup(html, 'html.parser')
    return soup.get_text(separator='\n', strip=True)

def extract_text_from_json(path):
    with open(path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    # JSON全体をインデント付き文字列に変換して返す
    return json.dumps(data, ensure_ascii=False, indent=2)

def extract_text_from_xlsx(path):
    wb = load_workbook(path, data_only=True)
    text = ''
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = [str(cell) if cell is not None else '' for cell in row]
            text += '\t'.join(row_text) + '\n'
    return text

def extract_text_from_pptx(path):
    prs = Presentation(path)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

def get_info(context):
    system_context = """あなたは優秀なアシスタントです。社内にあるドキュメントの内容を読み解き、わかりやすく要約し、キーワードを抽出します。
ナレッジベースを作成して、RAGに活用していきます。以下の制約条件と形式を守って、JSON形式で出力してください。
###制約条件
- 与えられるコンテキストは、ドキュメントをチャンクした文章です。与えられたチャンクの部分を要約し、summaryの値として出力します。要約した内容には、重要なキーワードは含めるようにしてください。
- 与えられたチャンクの文章に対して1文でタイトルを付与します。titleの値として出力します。
- 本チャンク内で検索に活用する重要なキーワードを抽出する。キーワードは25個以内とします
- 出力形式を守ります
###出力形式
summary: <チャンクした部分を要約した内容>
title: <チャンクした部分のタイトル>
Keywords: ["keyword1", "Keyword2", ...]  """

    user_request = "以下のコンテキストから制約条件と出力形式を必ず守って、JSON形式で出力をしてください。最初から最後まで注意深く読み込んでください。\n最高の仕事をしましょう。あなたならできる！\n###コンテキスト" + str(context)

    messages = [
        {"role": "system", "content": system_context},
        {"role": "user", "content": user_request}
    ]

    response = client.chat.completions.create(
        model=AZURE_OPENAI_CHAT_MODEL,
        messages=messages,
        temperature=0.0,
        max_tokens=AZURE_OPENAI_CHAT_MAX_TOKENS,
        response_format={"type": "json_object"},
    )

    content = json.loads(response.choices[0].message.content)
    return {
        'title': content['title'],
        'summary': content['summary'],
        'Keywords': content['Keywords']
    }

def get_vector(content):
    resp = client.embeddings.create(model=AZURE_OPENAI_EMBED_MODEL, input=content)
    return resp.data[0].embedding

def check_index_exists(name):
    try:
        index_client.get_index(name)
        return True
    except ResourceNotFoundError:
        return False

def delete_index(name):
    index_client.delete_index(name)

def create_index(name, json_file_path):
    with open(json_file_path, "r", encoding='utf-8') as f:
        data = json.load(f)
    data["name"] = name
    resp = requests.post(
        f"{AI_SEARCH_ENDPOINT}/indexes?api-version={AI_SEARCH_API_VERSION}",
        data=json.dumps(data),
        headers={"Content-Type": "application/json", "api-key": AI_SEARCH_KEY},
    )
    if not str(resp.status_code).startswith("2"):
        raise Exception(resp.text)
    return resp.status_code

def add_documents(index_name, docs):
    search_client = SearchClient(
        endpoint=AI_SEARCH_ENDPOINT,
        credential=credential,
        index_name=index_name,
    )
    search_client.upload_documents(documents=docs)

def add_to_cosmos(item):
    container.upsert_item(item)

def process_file(file_path, index_name=None, args=None):
    print("process file:", file_path)
    if index_name is None:
        index_name = os.getenv("AI_SEARCH_INDEX_NAME")

    print("use_blob_storage:", USE_BLOB_STORAGE)
    download_file_path = file_path
    metadata = {}
    if USE_BLOB_STORAGE and args.blob:
        blob_client = blob_container_client.get_blob_client(os.path.basename(file_path))
        download_file_path = os.path.join(os.getcwd(), os.path.basename(file_path))
        with open(download_file_path, "wb") as download_file:
            download_data = blob_client.download_blob()
            download_data.readinto(download_file)
        print(f"Downloaded {file_path} from Blob Storage to {download_file_path}.")
        metadata = blob_client.get_blob_properties().metadata

    ext = os.path.splitext(download_file_path)[1].lower()
    if ext == ".pdf":
        content = get_content_from_document(download_file_path)
    elif ext == ".txt":
        with open(download_file_path, 'r', encoding='utf-8') as f:
            content = f.read()
    elif ext == ".docx":
        content = extract_text_from_docx(download_file_path)
    elif ext == ".xlsx":
        content = extract_text_from_xlsx(download_file_path)
    elif ext == ".pptx":
        content = extract_text_from_pptx(download_file_path)
    elif ext == ".png":
        content = get_content_from_document(download_file_path)
    elif ext in [".html", ".htm"]:
        content = extract_text_from_html(download_file_path)
    elif ext == ".json":
        content = extract_text_from_json(download_file_path)
    else:
        print("unsupported file format:", download_file_path)
        return

    chunks = chunk_content(content, max_chunk_token_size, overlap_token_rate, overlap_type)

    if USE_BLOB_STORAGE:
        os.remove(download_file_path)

    file_name = os.path.basename(file_path)
    decoded_metadata = {k: urllib.parse.unquote(v) for k, v in metadata.items()}
    index_docs = []
    for chunk_no, chunk in enumerate(chunks):
        print("enrichment chunk:", f"{chunk_no+1}/{len(chunks)}")
        docinfo = get_info(chunk)
        id_base = f"{file_name}_{chunk_no}"
        id_hash = hashlib.sha256(id_base.encode('utf-8')).hexdigest()
        index_doc = {
            "id": id_hash,
            "fileName": file_name,
            "chunkNo": chunk_no,
            "content": chunk,
            "title": docinfo['title'],
            "summary": docinfo['summary'],
            "keywords": docinfo['Keywords'],
            "contentVector": get_vector(docinfo['summary']),
            "PageName": decoded_metadata.get("PageName", ""),
            "NoteName": decoded_metadata.get("NoteName", ""),
            "SectionName": decoded_metadata.get("SectionName", ""),
            "webURL": metadata.get("webURL", "")
        }
        index_docs.append(index_doc)
        add_to_cosmos(index_doc)

    if not check_index_exists(index_name):
        print("create index:", index_name)
        create_index(index_name, "index.json")

    print("upload documents to index:", index_name)
    add_documents(index_name, index_docs)

def main():
    parser = argparse.ArgumentParser(description='Process files for RAG application.')
    parser.add_argument('--file', type=str, help='Path to the file to process.')
    parser.add_argument('--dir', type=str, help='Path to the directory containing files to process.')
    parser.add_argument('--blob', action='store_true', help='Process all files in Blob Storage.')
    parser.add_argument('--index', type=str, help='Name of the Azure Cognitive Search index.')

    args = parser.parse_args()
    index_name = args.index if args.index else os.getenv("AI_SEARCH_INDEX_NAME")

    if args.file:
        process_file(args.file, index_name, args)
    elif args.dir:
        for root, dirs, files in os.walk(args.dir):
            for file in files:
                file_path = os.path.join(root, file)
                process_file(file_path, index_name, args)
    elif args.blob:
        print("Processing all files in Blob Storage...")
        blob_list = blob_container_client.list_blobs()
        for blob in blob_list:
            process_file(blob.name, index_name, args)
    else:
        print("Please specify a file, directory, or use --blob to process files from Blob Storage.")
        parser.print_help()

if __name__ == "__main__":
    main()
